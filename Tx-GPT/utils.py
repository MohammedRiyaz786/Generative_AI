import pdfplumber
import csv
import pandas as pd
from langchain_core.documents import Document
from io import StringIO
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_pdf_text(pdf_docs):
    """Extract text from PDFs, including tables."""
    text = ""
    documents = []
    
    for pdf in pdf_docs:
        with pdfplumber.open(pdf) as pdf_reader:
            for page_num, page in enumerate(pdf_reader.pages):
                # Extracting page text
                page_text = page.extract_text(x_tolerance=3, y_tolerance=3) or ""
                text += page_text + "\n"
                
                # Extracting tables
                tables = page.extract_tables()
                for table in tables:
                    table_text = "Table:\n"
                    for row in table:
                        filtered_row = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                        if filtered_row:
                            table_text += " | ".join(filtered_row) + "\n"
                    text += table_text + "\n"
                    documents.append(Document(page_content=table_text, metadata={'source': 'table', 'page': page_num + 1}))
                
                # Add non-table text as a separate document
                documents.append(Document(page_content=page_text, metadata={'source': 'pdf_text', 'page': page_num + 1}))
                    
    return text, documents

def get_non_table_pdf_text(pdf_docs):
    """Extract text from non-table PDFs with improved handling."""
    text = ""
    documents = []
    
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page_num, page in enumerate(pdf_reader.pages):
            # Extracting page text
            page_text = page.extract_text() or ""
            
            # Improved text cleaning
            cleaned_text = ' '.join(page_text.split())  # Remove extra whitespace
            cleaned_text = cleaned_text.replace('-\n', '')  # Handle hyphenated words
            
            # Add page number to metadata
            metadata = {'source': 'pdf_non_table', 'page': page_num + 1}
            
            # Create smaller chunks of text
            chunks = [cleaned_text[i:i+1000] for i in range(0, len(cleaned_text), 800)]
            
            for chunk in chunks:
                text += chunk + "\n"
                documents.append(Document(page_content=chunk, metadata=metadata))
                
    return text, documents

def get_csv_text(csv_file):
    """Extract text from CSV files."""
    text = ""
    csv_file.seek(0)
    content = csv_file.read().decode('utf-8')
    csv_reader = csv.reader(StringIO(content))
    for row in csv_reader:
        filtered_row = [cell.strip() for cell in row if cell.strip()]
        if filtered_row:
            text += " | ".join(filtered_row) + "\n"
    return text

def get_excel_text(excel_files):
    """Extract text from Excel files."""
    text = ""
    documents = []
    for excel_file in excel_files:
        excel_data = pd.read_excel(excel_file, sheet_name=None)
        for sheet_name, sheet_data in excel_data.items():
            sheet_text = f"Sheet: {sheet_name}\n"
            for _, row in sheet_data.iterrows():
                filtered_row = [str(cell).strip() for cell in row if pd.notna(cell) and str(cell).strip()]
                if filtered_row:
                    sheet_text += " | ".join(filtered_row) + "\n"
            text += sheet_text + "\n"
            documents.append(Document(page_content=sheet_text, metadata={'source': f"Excel Sheet {sheet_name}"}))
    return text, documents



def get_ppt_text(ppt_files):
    """Extract text from PowerPoint files, including tables and complex shapes."""
    text = ""
    documents = []
    
    for ppt_file in ppt_files:
        prs = Presentation(ppt_file)
        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"Slide {slide_num + 1}:\n"
            
            # Extracting text from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text += shape.text + "\n"
            
            # Extracting text from tables
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_text = f"Slide {slide_num + 1}: Table\n"
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        table_text += row_text + "\n"
                    text += table_text + "\n"
                    documents.append(Document(page_content=table_text, metadata={'source': f"PowerPoint Slide {slide_num + 1}"}))

            text += slide_text + "\n"
            documents.append(Document(page_content=slide_text, metadata={'source': f"PowerPoint Slide {slide_num + 1}"}))
    
    return text, documents


def get_word_text(word_files):
    """Extract text from Word files."""
    text = ""
    documents = []
    for word_file in word_files:
        doc = DocxDocument(word_file)
        for para_num, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text += para.text + "\n"
                documents.append(Document(page_content=para.text, metadata={'source': f"Word Paragraph {para_num + 1}"}))
    return text, documents